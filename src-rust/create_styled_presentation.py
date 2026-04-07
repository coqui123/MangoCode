
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# --- 1. DESIGN SYSTEM ---

# a. Color Palette (Inspired by Mangoes)
COLORS = {
    "bg_light": RGBColor(250, 248, 245),
    "bg_dark": RGBColor(45, 50, 65),
    "text_light": RGBColor(250, 248, 245),
    "text_dark": RGBColor(25, 28, 35),
    "text_accent": RGBColor(255, 180, 0), # Mango yellow/orange
    "accent1": RGBColor(255, 180, 0), # Mango yellow/orange
    "accent2": RGBColor(104, 182, 137), # Leafy green
    "accent3": RGBColor(218, 80, 50), # Ripe mango red
    "card_bg_light": RGBColor(255, 255, 255),
    "card_bg_dark": RGBColor(60, 65, 80),
}

# b. Typography
FONT_FAMILY = "Calibri" # A widely available font
TYPE_SCALE = {
    "hero": {"size": Pt(48), "bold": True},
    "title": {"size": Pt(32), "bold": True},
    "subtitle": {"size": Pt(20), "bold": False},
    "body": {"size": Pt(14), "bold": False},
    "label": {"size": Pt(12), "bold": True},
    "caption": {"size": Pt(10), "bold": False},
}

# c. Spacing & Layout
SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5
MARGIN = Inches(0.5)
TITLE_POS_Y = Inches(0.75)
CONTENT_START_Y = Inches(1.8)

# --- 2. HELPER/COMPONENT FUNCTIONS ---

def create_presentation():
    """Creates a new presentation with standard slide dimensions."""
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    return prs

def add_slide(prs, bg_mode='light'):
    """Adds a new blank slide and sets its background color."""
    slide_layout = prs.slide_layouts[6] # BLANK layout
    slide = prs.slides.add_slide(slide_layout)
    bg_color = COLORS['bg_light'] if bg_mode == 'light' else COLORS['bg_dark']
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color
    return slide

def add_text(slide, text, pos, size_info, color, alignment=pptx.enum.text.PP_ALIGN.LEFT):
    """Adds a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(pos[0], pos[1], pos[2], pos[3])
    tf = txBox.text_frame
    tf.clear() 
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_FAMILY
    p.font.size = size_info["size"]
    p.font.bold = size_info["bold"]
    p.font.color.rgb = color
    p.alignment = alignment
    tf.vertical_anchor = pptx.enum.text.MSO_ANCHOR.TOP
    tf.word_wrap = True
    return txBox

def add_slide_title(slide, title, bg_mode='light'):
    """Adds a formatted title to the slide."""
    text_color = COLORS['text_dark'] if bg_mode == 'light' else COLORS['text_light']
    pos = (MARGIN, TITLE_POS_Y, Inches(SLIDE_WIDTH_INCHES) - 2 * MARGIN, Inches(1))
    return add_text(slide, title, pos, TYPE_SCALE["title"], text_color)

def add_accent_bar(slide, bg_mode='light'):
    """Adds a small decorative accent bar to the top of the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        MARGIN, Inches(0.3), Inches(2.5), Inches(0.08)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['accent1']
    shape.line.fill.background()

def add_card(slide, pos, bg_mode='light'):
    """Adds a styled card shape."""
    card_bg = COLORS['card_bg_light'] if bg_mode == 'light' else COLORS['card_bg_dark']
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos[0], pos[1], pos[2], pos[3])
    shape.shadow.inherit = False
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = card_bg
    shape.line.fill.background() # No outline
    return shape

# --- 3. SLIDE GENERATION ---

def generate_mango_presentation():
    """Main function to generate the full presentation."""
    prs = create_presentation()

    # --- Slide 1: Title Slide ---
    slide1 = add_slide(prs, 'dark')
    add_accent_bar(slide1, 'dark')
    add_text(slide1, "MANGO", 
             (MARGIN, Inches(2.5), Inches(8), Inches(2)), 
             TYPE_SCALE["hero"], COLORS['accent1'])
    add_text(slide1, "The Undisputed King of Fruits", 
             (MARGIN, Inches(3.8), Inches(8), Inches(1)), 
             TYPE_SCALE["subtitle"], COLORS['text_light'])
    add_text(slide1, "An Educational Primer", 
             (MARGIN, Inches(6.5), Inches(5), Inches(0.5)), 
             TYPE_SCALE["label"], COLORS['accent2'])

    # --- Slide 2: Agenda ---
    slide2 = add_slide(prs, 'light')
    add_accent_bar(slide2)
    add_slide_title(slide2, "What we'll cover")
    agend-items = [
        "A Brief History",
        "Popular Varieties",
        "Nutritional Power",
        "How to Prepare & Eat",
        "Fun Facts"
    ]
    for i, item in enumerate(agend-items):
        y_pos = CONTENT_START_Y + Inches(i * 0.7)
        add_card(slide2, (MARGIN, y_pos, Inches(6), Inches(0.6)))
        add_text(slide2, f"0{i+1}", 
                 (MARGIN + Inches(0.2), y_pos + Inches(0.05), Inches(0.5), Inches(0.5)), 
                 TYPE_SCALE["title"], COLORS['accent1'])
        add_text(slide2, item, 
                 (MARGIN + Inches(1), y_pos + Inches(0.15), Inches(5), Inches(0.5)), 
                 TYPE_SCALE["subtitle"], COLORS['text_dark'])

    # --- Slide 3: History ---
    slide3 = add_slide(prs, 'dark')
    add_accent_bar(slide3, 'dark')
    add_slide_title(slide3, "From Ancient Royalty to Global Star", 'dark')
    add_text(slide3, "The mango originated in South Asia over 4,000 years ago, where it was considered a sacred fruit. It was a symbol of love and friendship, often exchanged as a gift by royalty.",
             (MARGIN, CONTENT_START_Y, Inches(6), Inches(3)),
             TYPE_SCALE["body"], COLORS['text_light'])
    add_text(slide3, "4,000+", 
             (Inches(8), Inches(2.5), Inches(4), Inches(2)),
             TYPE_SCALE["hero"], COLORS['accent1'])
    add_text(slide3, "YEARS OF HISTORY", 
             (Inches(8), Inches(3.8), Inches(4), Inches(1)),
             TYPE_SCALE["label"], COLORS['text_light'])
             
    # --- Slide 4: Varieties (Placeholder) ---
    slide4 = add_slide(prs, 'light')
    add_accent_bar(slide4)
    add_slide_title(slide4, "Hundreds of Varieties, Unique Flavors")
    # This is a placeholder for where images and descriptions of mangoes would go
    varieties = ["Tommy Atkins", "Honey (Ataulfo)", "Kent", "Keitt"]
    for i, variety in enumerate(varieties):
        x_pos = MARGIN + Inches(i * 3.2)
        add_card(slide4, (x_pos, CONTENT_START_Y, Inches(3), Inches(4)))
        add_text(slide4, variety, (x_pos, CONTENT_START_Y + Inches(0.2), Inches(3), Inches(0.5)), TYPE_SCALE["subtitle"], COLORS['text_dark'], alignment=pptx.enum.text.PP_ALIGN.CENTER)
        add_text(slide4, "[Image Placeholder]", (x_pos, CONTENT_START_Y + Inches(1.5), Inches(3), Inches(0.5)), TYPE_SCALE["body"], COLORS['text_dark'], alignment=pptx.enum.text.PP_ALIGN.CENTER)


    # --- Slide 5: Nutrition ---
    slide5 = add_slide(prs, 'dark')
    add_accent_bar(slide5, 'dark')
    add_slide_title(slide5, "A Nutritional Powerhouse", 'dark')
    add_text(slide5, "Mangos aren't just delicious; they are packed with over 20 vitamins and minerals.",
             (MARGIN, CONTENT_START_Y, Inches(SLIDE_WIDTH_INCHES - 2*MARGIN), Inches(1)),
             TYPE_SCALE["body"], COLORS['text_light'])
    # ... more slides would be generated here ...
    
    # --- Last Slide: Closing ---
    slide_final = add_slide(prs, 'light')
    add_accent_bar(slide_final)
    add_text(slide_final, "Go Enjoy a Mango!", 
             (MARGIN, Inches(3), Inches(12), Inches(1.5)), 
             TYPE_SCALE["hero"], COLORS['text_dark'], alignment=pptx.enum.text.PP_ALIGN.CENTER)

    return prs

# --- 4. SAVE PRESENTATION ---
if __name__ == '__main__':
    presentation = generate_mango_presentation()
    presentation.save('mango_presentation.pptx')
    print("Presentation 'mango_presentation.pptx' generated successfully.")

